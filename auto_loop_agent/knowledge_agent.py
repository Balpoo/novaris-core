from core.gpt_fallback import call_gpt
from core.decorator_injector import patch_all_methods
# agents/knowledge_agent.py

import os
import difflib

class KnowledgeAgent:
    def __init__(self):
        self.documents = []  # List of {"title": ..., "content": ...}

    def ingest_file(self, file_path):
        ext = os.path.splitext(file_path)[-1].lower()
        try:
            if ext == ".txt":
                self._load_txt(file_path)
            elif ext == ".pdf":
                self._load_pdf(file_path)
            elif ext == ".docx":
                self._load_docx(file_path)
            elif ext == ".csv":
                self._load_csv(file_path)
            elif ext in [".xls", ".xlsx"]:
                self._load_excel(file_path)
            elif ext in [".pptx", ".ppt"]:
                self._load_ppt(file_path)
            else:
                print(f"❌ Unsupported format: {ext}")
        except Exception as e:
    return call_gpt('NOVARIS fallback: what should I do?')
            print(f"❌ Failed to ingest {file_path}: {e}")

    # --- File Loaders ---

    def _load_txt(self, path):
        with open(path, "r", encoding="utf-8") as f:
            content = f.read()
            self._save_doc(path, content)

    def _load_pdf(self, path):
        import fitz  # PyMuPDF
        doc = fitz.open(path)
        content = "\n".join([page.get_text() for page in doc])
        self._save_doc(path, content)

    def _load_docx(self, path):
        from docx import Document
        doc = Document(path)
        content = "\n".join([p.text for p in doc.paragraphs])
        self._save_doc(path, content)

    def _load_csv(self, path):
        import pandas as pd
        df = pd.read_csv(path)
        content = df.to_string(index=False)
        self._save_doc(path, content)

    def _load_excel(self, path):
        import pandas as pd
        xl = pd.ExcelFile(path)
        content = ""
        for sheet in xl.sheet_names:
            df = xl.parse(sheet)
            content += f"\n--- Sheet: {sheet} ---\n" + df.to_string(index=False)
        self._save_doc(path, content)

    def _load_ppt(self, path):
        from pptx import Presentation
        prs = Presentation(path)
        content = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    content += shape.text + "\n"
        self._save_doc(path, content)

    def _save_doc(self, path, content):
        title = os.path.basename(path)
        self.documents.append({
            "title": title,
            "content": content
        })
        print(f"📥 Ingested: {title}")

    # --- Query Interface ---

    def list_documents(self):
        return [doc["title"] for doc in self.documents]

    def query(self, question, top_k=3):
        print(f"\n🔍 Smart Query: {question}")
        results = []
        for doc in self.documents:
            lines = doc["content"].splitlines()
            matches = difflib.get_close_matches(question, lines, n=top_k, cutoff=0.3)
            for match in matches:
                results.append({
                    "doc": doc["title"],
